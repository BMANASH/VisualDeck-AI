import streamlit as st
from PIL import Image
import io
import json
import math
import time
from pypdf import PdfReader
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import pandas as pd

# =========================================================
# 1. Page Configuration
# =========================================================
st.set_page_config(
    page_title="VisualDeck AI | Intelligent Presentation Engine",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded"
)

# =========================================================
# 2. Futuristic Styling (CSS) + Glassmorphic Loading Overlay
# =========================================================
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@400;500;600;700&family=JetBrains+Mono:wght@400;500;600&display=swap');

    html, body, [class*="css"]  {
        font-family: 'Space Grotesk', sans-serif;
    }

    .stApp {
        background:
            radial-gradient(circle at 15% 10%, rgba(0, 210, 255, 0.10) 0%, transparent 45%),
            radial-gradient(circle at 85% 0%, rgba(157, 0, 255, 0.10) 0%, transparent 40%),
            radial-gradient(circle at 50% 100%, rgba(0, 255, 178, 0.07) 0%, transparent 45%),
            #05070C;
        background-attachment: fixed;
    }

    /* Animated grid backdrop behind hero */
    .grid-backdrop {
        position: relative;
        padding: 4px 0 10px 0;
    }
    .grid-backdrop::before {
        content: "";
        position: absolute;
        top: -40px; left: -60px; right: -60px; height: 220px;
        background-image:
            linear-gradient(rgba(0, 210, 255, 0.07) 1px, transparent 1px),
            linear-gradient(90deg, rgba(0, 210, 255, 0.07) 1px, transparent 1px);
        background-size: 34px 34px;
        mask-image: linear-gradient(to bottom, black 20%, transparent 90%);
        pointer-events: none;
        z-index: 0;
    }

    .kpi-card {
        position: relative;
        background: linear-gradient(160deg, rgba(255,255,255,0.05), rgba(255,255,255,0.015));
        border: 1px solid rgba(0, 210, 255, 0.20);
        border-radius: 14px;
        padding: 16px 18px;
        backdrop-filter: blur(14px);
        box-shadow: 0 4px 24px 0 rgba(0,0,0,0.45), inset 0 1px 0 rgba(255,255,255,0.04);
        margin-bottom: 10px;
        overflow: hidden;
        transition: border-color 0.2s ease;
    }
    .kpi-card:hover { border-color: rgba(0, 210, 255, 0.55); }
    .kpi-card::after {
        content: "";
        position: absolute; top: 0; left: 0; right: 0; height: 2px;
        background: linear-gradient(90deg, #00D2FF, #9D00FF, transparent);
        opacity: 0.8;
    }
    .kpi-title {
        font-family: 'JetBrains Mono', monospace;
        font-size: 0.66rem;
        text-transform: uppercase;
        letter-spacing: 1.4px;
        color: #7C8CA3;
        margin-bottom: 6px;
        font-weight: 500;
    }
    .kpi-value {
        font-size: 1.35rem;
        font-weight: 700;
        color: #00D2FF;
        letter-spacing: -0.5px;
    }
    .kpi-desc {
        font-size: 0.74rem;
        color: #5C6B80;
        margin-top: 3px;
    }
    .badge-pill {
        display: inline-flex;
        align-items: center;
        gap: 6px;
        padding: 5px 14px;
        border-radius: 50px;
        font-family: 'JetBrains Mono', monospace;
        font-size: 0.68rem;
        font-weight: 600;
        letter-spacing: 0.8px;
        background: rgba(0, 210, 255, 0.08);
        border: 1px solid rgba(0, 210, 255, 0.35);
        color: #00D2FF;
        margin-bottom: 10px;
    }
    .badge-pill .dot {
        width: 6px; height: 6px; border-radius: 50%;
        background: #00FFB2;
        box-shadow: 0 0 8px 2px rgba(0,255,178,0.7);
        display: inline-block;
        animation: pulse-dot 1.6s infinite ease-in-out;
    }
    @keyframes pulse-dot {
        0%, 100% { opacity: 1; transform: scale(1); }
        50% { opacity: 0.4; transform: scale(0.7); }
    }
    .sidebar-kpi {
        background: rgba(14, 20, 32, 0.75);
        border: 1px solid rgba(0, 210, 255, 0.18);
        border-radius: 12px;
        padding: 12px 14px;
        margin-bottom: 10px;
        backdrop-filter: blur(10px);
    }
    .sidebar-kpi-title {
        font-family: 'JetBrains Mono', monospace;
        font-size: 0.62rem;
        color: #7C8CA3;
        font-weight: 600;
        text-transform: uppercase;
        letter-spacing: 1px;
    }
    .sidebar-kpi-val {
        font-size: 0.9rem;
        color: #F1F5F9;
        font-weight: 600;
        margin-top: 3px;
    }
    .hero-title {
        position: relative;
        z-index: 1;
        font-size: 2.6rem;
        font-weight: 700;
        background: linear-gradient(100deg, #FFFFFF 10%, #8FE9FF 55%, #B18CFF 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 2px;
        letter-spacing: -1px;
    }
    .hero-subtitle {
        position: relative;
        z-index: 1;
        color: #8CA0B8;
        font-size: 1rem;
        margin-bottom: 20px;
        max-width: 640px;
    }
    .section-label {
        font-family: 'JetBrains Mono', monospace;
        font-size: 0.72rem;
        color: #00D2FF;
        letter-spacing: 1.5px;
        text-transform: uppercase;
        margin: 6px 0 2px 0;
        opacity: 0.85;
    }

    /* ---- Glassmorphic Loading Overlay ---- */
    .glass-loader-backdrop {
        position: fixed;
        inset: 0;
        background: rgba(3, 6, 12, 0.55);
        backdrop-filter: blur(3px);
        z-index: 9998;
        display: flex;
        align-items: center;
        justify-content: center;
    }
    .glass-loader-box {
        width: 320px;
        padding: 28px 26px 24px 26px;
        border-radius: 20px;
        background: linear-gradient(160deg, rgba(255,255,255,0.09), rgba(255,255,255,0.02));
        border: 1px solid rgba(0, 210, 255, 0.35);
        box-shadow: 0 20px 60px rgba(0,0,0,0.6), inset 0 1px 0 rgba(255,255,255,0.08);
        backdrop-filter: blur(20px);
        text-align: center;
        z-index: 9999;
        font-family: 'Space Grotesk', sans-serif;
    }
    .glass-loader-ring {
        width: 54px; height: 54px;
        margin: 0 auto 16px auto;
        border-radius: 50%;
        border: 3px solid rgba(0, 210, 255, 0.15);
        border-top-color: #00D2FF;
        border-right-color: #9D00FF;
        animation: spin 0.9s linear infinite;
    }
    @keyframes spin { to { transform: rotate(360deg); } }
    .glass-loader-title {
        color: #F1F5F9;
        font-size: 0.98rem;
        font-weight: 600;
        margin-bottom: 4px;
    }
    .glass-loader-phase {
        color: #00D2FF;
        font-family: 'JetBrains Mono', monospace;
        font-size: 0.72rem;
        letter-spacing: 0.5px;
    }
    .glass-loader-bar-track {
        margin-top: 14px;
        width: 100%;
        height: 4px;
        border-radius: 4px;
        background: rgba(255,255,255,0.08);
        overflow: hidden;
    }
    .glass-loader-bar-fill {
        height: 100%;
        width: 40%;
        border-radius: 4px;
        background: linear-gradient(90deg, #00D2FF, #9D00FF);
        animation: loadbar 1.3s ease-in-out infinite;
    }
    @keyframes loadbar {
        0% { transform: translateX(-100%); }
        100% { transform: translateX(250%); }
    }
</style>
""", unsafe_allow_html=True)


def render_loader(placeholder, title, phase_text):
    """Render a glassmorphic loading popup into a placeholder. CSS animation keeps
    spinning in the browser even while the Python thread is blocked on the API call."""
    placeholder.markdown(f"""
    <div class="glass-loader-backdrop">
        <div class="glass-loader-box">
            <div class="glass-loader-ring"></div>
            <div class="glass-loader-title">{title}</div>
            <div class="glass-loader-phase">{phase_text}</div>
            <div class="glass-loader-bar-track"><div class="glass-loader-bar-fill"></div></div>
        </div>
    </div>
    """, unsafe_allow_html=True)


# =========================================================
# 3. Helper Functions: Colors, Text Measurement, Images
# =========================================================
def hex_to_rgb(hex_code, default_hex="00D2FF"):
    try:
        clean_hex = str(hex_code).lstrip('#')
        if len(clean_hex) != 6:
            clean_hex = default_hex.lstrip('#')
        return RGBColor(*(int(clean_hex[i:i + 2], 16) for i in (0, 2, 4)))
    except Exception:
        return RGBColor(0, 210, 255)


SAFE_FONTS = {"Calibri", "Arial", "Segoe UI", "Georgia", "Verdana",
              "Trebuchet MS", "Century Gothic", "Garamond", "Tahoma"}


def safe_font(name, fallback="Calibri"):
    if isinstance(name, str) and name.strip() in SAFE_FONTS:
        return name.strip()
    return fallback


def estimate_block_height_in(text, font_size_pt, box_width_in, line_spacing=1.22, min_lines=1, bold=False):
    """Rough text-wrapping estimate (no real font metrics available at build time).
    Used to size containers to their actual content instead of using fixed heights."""
    if not text:
        text = ""
    avg_char_width_in = (font_size_pt * (0.62 if bold else 0.52)) / 72.0
    chars_per_line = max(1, int(box_width_in / avg_char_width_in))
    lines = max(min_lines, math.ceil(len(str(text)) / chars_per_line)) if text else min_lines
    line_height_in = (font_size_pt * line_spacing) / 72.0
    return lines * line_height_in


def fit_font_size_to_width(text, max_width_in, max_size=28, min_size=15, bold=True):
    """Shrink font size (like PowerPoint's own autofit) until the text is estimated to fit
    on one line within max_width_in. Prevents large KPI numbers/words from overflowing their card.
    Deliberately conservative (overestimates glyph width) since wide fonts like Century Gothic
    can otherwise wrap where the estimate predicted a single line."""
    text = str(text) if text else ""
    if not text:
        return max_size
    factor = 0.74 if bold else 0.55
    size = max_size
    while size > min_size:
        est_width = len(text) * (size * factor) / 72.0
        if est_width <= max_width_in:
            break
        size -= 1
    return size


def estimate_lines_at_size(text, font_size_pt, box_width_in, bold=True):
    text = str(text) if text else ""
    if not text:
        return 1
    factor = 0.74 if bold else 0.55
    est_width = len(text) * (font_size_pt * factor) / 72.0
    return max(1, math.ceil(est_width / box_width_in))


def fit_dimensions(img_w, img_h, max_w_in, max_h_in):
    """Preserve aspect ratio while fitting inside a bounding box."""
    ratio = min(max_w_in / img_w, max_h_in / img_h)
    return img_w * ratio, img_h * ratio


def pil_to_stream(pil_img, fmt="PNG"):
    buf = io.BytesIO()
    if fmt == "JPEG" and pil_img.mode != "RGB":
        pil_img = pil_img.convert("RGB")
    pil_img.save(buf, format=fmt)
    buf.seek(0)
    return buf


def add_rounded_card(slide, left, top, width, height, fill_rgb, line_rgb, line_w_pt=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb
    shape.line.width = Pt(line_w_pt)
    shape.shadow.inherit = False
    return shape


def add_icon_badge(slide, left, top, size, icon_char, bg_rgb, fg_rgb):
    """Small circular badge with an AI-chosen unicode/emoji glyph — avoids hardcoded icon mapping."""
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg_rgb
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(icon_char) if icon_char else "•"
    p.font.size = Pt(max(10, int(size * 28)))
    p.font.color.rgb = fg_rgb
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return badge


# =========================================================
# 4. PowerPoint Build Engine (dynamic sizing + real images)
# =========================================================
def build_powerpoint(deck_data, source_images=None):
    source_images = source_images or []  # list of PIL Images, index-aligned with what Gemini saw

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    palette = deck_data.get("theme", {})
    bg_color = hex_to_rgb(palette.get("background_hex", "0A0E17"))
    primary_color = hex_to_rgb(palette.get("primary_accent_hex", "00D2FF"))
    secondary_color = hex_to_rgb(palette.get("secondary_accent_hex", palette.get("primary_accent_hex", "9D00FF")))
    card_color = hex_to_rgb(palette.get("card_bg_hex", "141C2E"))
    text_color = hex_to_rgb(palette.get("text_color_hex", "FFFFFF"))
    secondary_text = hex_to_rgb(palette.get("secondary_text_hex", "94A3B8"))
    heading_font = safe_font(palette.get("heading_font"))
    body_font = safe_font(palette.get("body_font"))

    MARGIN = 0.8
    CONTENT_W = 13.333 - 2 * MARGIN
    CONTENT_TOP = 1.9
    CONTENT_BOTTOM = 6.7  # leave room for footer

    slides_list = deck_data.get("slides", [])

    for slide_info in slides_list:
        slide = prs.slides.add_slide(blank_layout)
        layout = slide_info.get("layout_type", "cards")

        # --- Background ---
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = bg_color
        bg_shape.line.fill.background()
        bg_shape.shadow.inherit = False

        # --- Section-break slides get a completely different, simpler treatment ---
        if layout == "section_break":
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(2.68), Inches(1.1), Inches(0.09))
            accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = primary_color; accent_bar.line.fill.background()

            title_box = slide.shapes.add_textbox(Inches(MARGIN), Inches(2.95), Inches(CONTENT_W), Inches(1.3))
            tf = title_box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_info.get("title", "")
            p.font.size = Pt(38); p.font.bold = True; p.font.name = heading_font
            p.font.color.rgb = text_color

            if slide_info.get("subtitle"):
                sub_box = slide.shapes.add_textbox(Inches(MARGIN), Inches(3.75), Inches(CONTENT_W), Inches(0.8))
                stf = sub_box.text_frame; stf.word_wrap = True
                sp = stf.paragraphs[0]
                sp.text = slide_info.get("subtitle")
                sp.font.size = Pt(15); sp.font.name = body_font
                sp.font.color.rgb = secondary_text

            footer_box = slide.shapes.add_textbox(Inches(MARGIN), Inches(6.85), Inches(CONTENT_W), Inches(0.4))
            f_p = footer_box.text_frame.paragraphs[0]
            f_p.text = f"{deck_data.get('deck_title', 'Presentation')}  •  {slide_info.get('slide_number', '')}"
            f_p.font.size = Pt(9); f_p.font.color.rgb = secondary_text; f_p.font.name = body_font
            continue

        # --- Accent line ---
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(0.4), Inches(CONTENT_W), Inches(0.035))
        accent_line.fill.solid(); accent_line.fill.fore_color.rgb = primary_color; accent_line.line.fill.background()

        # --- Title & subtitle ---
        title_box = slide.shapes.add_textbox(Inches(MARGIN), Inches(0.55), Inches(CONTENT_W), Inches(1.1))
        tf = title_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_info.get("title", "")
        p.font.size = Pt(24); p.font.bold = True; p.font.name = heading_font
        p.font.color.rgb = primary_color

        if slide_info.get("subtitle"):
            p_sub = tf.add_paragraph()
            p_sub.text = slide_info.get("subtitle")
            p_sub.font.size = Pt(12); p_sub.font.name = body_font
            p_sub.font.color.rgb = secondary_text

        img_idx = slide_info.get("image_index", None)
        pil_img = None
        if img_idx is not None and isinstance(img_idx, int) and 0 <= img_idx < len(source_images):
            pil_img = source_images[img_idx]

        # =====================================================
        # LAYOUT: image_left / image_right — photo + detail panel
        # =====================================================
        if layout in ("image_left", "image_right") and pil_img is not None:
            img_box_w, img_box_h = 5.4, CONTENT_BOTTOM - CONTENT_TOP
            iw, ih = fit_dimensions(pil_img.width, pil_img.height, img_box_w, img_box_h)
            img_left = MARGIN if layout == "image_left" else (13.333 - MARGIN - iw)
            img_top = CONTENT_TOP + (img_box_h - ih) / 2

            frame = add_rounded_card(slide, img_left - 0.08, img_top - 0.08, iw + 0.16, ih + 0.16, card_color, primary_color, 1.2)
            stream = pil_to_stream(pil_img, "JPEG")
            slide.shapes.add_picture(stream, Inches(img_left), Inches(img_top), width=Inches(iw), height=Inches(ih))

            text_left = (MARGIN + img_box_w + 0.5) if layout == "image_left" else MARGIN
            text_w = CONTENT_W - img_box_w - 0.5
            cursor_top = CONTENT_TOP
            for bullet in slide_info.get("bullets", [])[:6]:
                icon = bullet.get("icon", "•") if isinstance(bullet, dict) else "•"
                b_text = bullet.get("text", str(bullet)) if isinstance(bullet, dict) else str(bullet)
                b_height = max(0.55, estimate_block_height_in(b_text, 13, text_w - 0.5) + 0.18)
                add_icon_badge(slide, text_left, cursor_top + 0.05, 0.32, icon, primary_color, bg_color)
                b_box = slide.shapes.add_textbox(Inches(text_left + 0.48), Inches(cursor_top), Inches(text_w - 0.48), Inches(b_height))
                btf = b_box.text_frame; btf.word_wrap = True
                bp = btf.paragraphs[0]
                bp.text = b_text
                bp.font.size = Pt(13); bp.font.name = body_font; bp.font.color.rgb = text_color
                cursor_top += b_height + 0.22

        # =====================================================
        # LAYOUT: image_full — full-bleed image with title overlay
        # =====================================================
        elif layout == "image_full" and pil_img is not None:
            has_caption = bool(slide_info.get("bullets"))
            cap_text = ""
            cap_h = 0.0
            if has_caption:
                cap_text = " · ".join([b.get("text", str(b)) if isinstance(b, dict) else str(b) for b in slide_info["bullets"][:3]])
                cap_h = min(estimate_block_height_in(cap_text, 12, CONTENT_W - 0.6) + 0.3, 1.1)

            reserved_for_caption = (cap_h + 0.2) if has_caption else 0.0
            image_area_h = (CONTENT_BOTTOM - CONTENT_TOP) - reserved_for_caption
            iw, ih = fit_dimensions(pil_img.width, pil_img.height, CONTENT_W, image_area_h)
            left = MARGIN + (CONTENT_W - iw) / 2
            top = CONTENT_TOP + (image_area_h - ih) / 2
            stream = pil_to_stream(pil_img, "JPEG")
            slide.shapes.add_picture(stream, Inches(left), Inches(top), width=Inches(iw), height=Inches(ih))

            if has_caption:
                cap_top = CONTENT_TOP + image_area_h + 0.2
                cap_box = add_rounded_card(slide, MARGIN, cap_top, CONTENT_W, cap_h, card_color, primary_color, 0.75)
                ctf = cap_box.text_frame; ctf.word_wrap = True
                ctf.margin_left = Inches(0.25); ctf.margin_top = Inches(0.12)
                ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
                cp = ctf.paragraphs[0]
                cp.text = cap_text
                cp.font.size = Pt(12); cp.font.name = body_font; cp.font.color.rgb = text_color

        # =====================================================
        # LAYOUT: KPI grid — dynamic card height, icon badges
        # =====================================================
        elif layout == "kpi_grid" and slide_info.get("kpis"):
            kpis = slide_info["kpis"]
            card_count = min(len(kpis), 4)
            gap = 0.35
            card_width = (CONTENT_W - (card_count - 1) * gap) / card_count
            text_w = card_width - 0.5  # available width inside each card for text

            # First pass: work out a font size per card that actually fits its value text,
            # and how many lines that value will need at that size (mirrors PowerPoint autofit).
            card_plans = []
            for kpi in kpis[:card_count]:
                value_text = str(kpi.get("value", ""))
                label_text = str(kpi.get("label", ""))
                value_font = fit_font_size_to_width(value_text, text_w, max_size=28, min_size=15, bold=True)
                value_lines = estimate_lines_at_size(value_text, value_font, text_w, bold=True)
                value_h = value_lines * (value_font * 1.18) / 72.0

                label_font = fit_font_size_to_width(label_text, text_w, max_size=12.5, min_size=10, bold=True)
                label_lines = estimate_lines_at_size(label_text, label_font, text_w, bold=True)
                label_h = label_lines * (label_font * 1.25) / 72.0

                desc_h = 0.0
                if kpi.get("desc"):
                    desc_h = estimate_block_height_in(kpi.get("desc"), 10, text_w) + 0.12

                # badge_row + value + small gap + label + desc + top/bottom padding
                total_h = 0.62 + value_h + 0.08 + label_h + desc_h + 0.35
                card_plans.append({
                    "value_font": value_font, "value_lines": value_lines,
                    "label_font": label_font, "desc_h": desc_h,
                    "height": min(max(total_h, 1.7), 3.4)
                })

            card_height = max(p["height"] for p in card_plans)

            for idx, kpi in enumerate(kpis[:card_count]):
                plan = card_plans[idx]
                left = MARGIN + idx * (card_width + gap)
                top = CONTENT_TOP
                add_rounded_card(slide, left, top, card_width, card_height, card_color, primary_color, 1.0)

                if kpi.get("icon"):
                    add_icon_badge(slide, left + 0.22, top + 0.22, 0.42, kpi.get("icon"), primary_color, bg_color)

                c_box = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.72), Inches(text_w), Inches(card_height - 0.85))
                c_tf = c_box.text_frame; c_tf.word_wrap = True
                p_val = c_tf.paragraphs[0]
                p_val.text = str(kpi.get("value", ""))
                p_val.font.size = Pt(plan["value_font"]); p_val.font.bold = True; p_val.font.name = heading_font
                p_val.font.color.rgb = primary_color

                p_lbl = c_tf.add_paragraph()
                p_lbl.text = str(kpi.get("label", ""))
                p_lbl.font.size = Pt(plan["label_font"]); p_lbl.font.bold = True; p_lbl.font.name = body_font
                p_lbl.font.color.rgb = text_color

                if kpi.get("desc"):
                    p_desc = c_tf.add_paragraph()
                    p_desc.text = str(kpi.get("desc"))
                    p_desc.font.size = Pt(10); p_desc.font.name = body_font
                    p_desc.font.color.rgb = secondary_text

        # =====================================================
        # LAYOUT: Table — header shading + zebra striping
        # =====================================================
        elif layout == "table" and slide_info.get("table"):
            table_data = slide_info["table"]
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])

            if headers and rows:
                num_rows = len(rows) + 1
                num_cols = len(headers)
                avail_h = min(4.6, 0.55 * num_rows + 0.2)
                t_shape = slide.shapes.add_table(num_rows, num_cols, Inches(MARGIN), Inches(CONTENT_TOP), Inches(CONTENT_W), Inches(avail_h))
                table = t_shape.table

                for c_idx, head in enumerate(headers):
                    cell = table.cell(0, c_idx)
                    cell.text = str(head)
                    cell.fill.solid(); cell.fill.fore_color.rgb = card_color
                    cell.margin_left = cell.margin_right = Inches(0.12)
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(12); para.font.bold = True
                        para.font.name = heading_font; para.font.color.rgb = primary_color

                for r_idx, row in enumerate(rows):
                    row_fill = bg_color if r_idx % 2 == 0 else card_color
                    for c_idx, val in enumerate(row):
                        cell = table.cell(r_idx + 1, c_idx)
                        cell.text = str(val)
                        cell.fill.solid(); cell.fill.fore_color.rgb = row_fill
                        cell.margin_left = cell.margin_right = Inches(0.12)
                        for para in cell.text_frame.paragraphs:
                            para.font.size = Pt(10.5); para.font.name = body_font
                            para.font.color.rgb = text_color

        # =====================================================
        # LAYOUT: Cards / bullets — dynamic height, icon badges
        # =====================================================
        else:
            bullets = slide_info.get("bullets", [])
            card_count = min(len(bullets), 6)
            if card_count > 0:
                gap = 0.2
                items = []
                for bullet in bullets[:card_count]:
                    b_text = bullet.get("text", str(bullet)) if isinstance(bullet, dict) else str(bullet)
                    icon = bullet.get("icon", "•") if isinstance(bullet, dict) else "•"
                    h = estimate_block_height_in(b_text, 13, CONTENT_W - 1.0) + 0.34
                    items.append((b_text, icon, max(h, 0.62)))

                total_h = sum(h for _, _, h in items) + gap * (card_count - 1)
                available = CONTENT_BOTTOM - CONTENT_TOP
                scale = min(1.0, available / total_h) if total_h > available else 1.0

                cursor_top = CONTENT_TOP
                for b_text, icon, h in items:
                    h_scaled = h * scale
                    add_rounded_card(slide, MARGIN, cursor_top, CONTENT_W, h_scaled, card_color, primary_color, 1.0)
                    add_icon_badge(slide, MARGIN + 0.22, cursor_top + h_scaled / 2 - 0.16, 0.32, icon, primary_color, bg_color)
                    b_box = slide.shapes.add_textbox(Inches(MARGIN + 0.7), Inches(cursor_top), Inches(CONTENT_W - 1.0), Inches(h_scaled))
                    b_tf = b_box.text_frame; b_tf.word_wrap = True
                    b_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    bp = b_tf.paragraphs[0]
                    bp.text = b_text
                    bp.font.size = Pt(13); bp.font.name = body_font; bp.font.color.rgb = text_color
                    cursor_top += h_scaled * scale + gap if scale < 1 else h_scaled + gap
                    cursor_top = cursor_top if scale >= 1 else cursor_top

        # --- Footer ---
        footer_box = slide.shapes.add_textbox(Inches(MARGIN), Inches(6.95), Inches(CONTENT_W), Inches(0.4))
        f_tf = footer_box.text_frame
        f_p = f_tf.paragraphs[0]
        f_p.text = f"{deck_data.get('deck_title', 'VisualDeck AI')}  •  Slide {slide_info.get('slide_number', '')}"
        f_p.font.size = Pt(9); f_p.font.color.rgb = secondary_text; f_p.font.name = body_font

    output_stream = io.BytesIO()
    prs.save(output_stream)
    output_stream.seek(0)
    return output_stream


# =========================================================
# 5. Dialog for Clean Full-Image Inspection Popups
# =========================================================
@st.dialog("🖼️ Source Image Preview", width="large")
def preview_modal(img_obj, filename):
    st.image(img_obj, caption=filename, use_container_width=True)


# =========================================================
# 6. Left Sidebar Overview
# =========================================================
with st.sidebar:
    st.markdown('<div class="badge-pill"><span class="dot"></span>PLATFORM CONTROL</div>', unsafe_allow_html=True)
    st.markdown("### ⚡ VisualDeck AI")
    st.caption("Autonomous visual-to-presentation synthesis engine.")
    st.markdown("---")

    st.markdown("""
    <div class="sidebar-kpi">
        <div class="sidebar-kpi-title">Core AI Architecture</div>
        <div class="sidebar-kpi-val" style="color:#00D2FF;">Gemini 3 Series</div>
    </div>
    <div class="sidebar-kpi">
        <div class="sidebar-kpi-title">Presentation Format</div>
        <div class="sidebar-kpi-val">16:9 Widescreen .PPTX</div>
    </div>
    <div class="sidebar-kpi">
        <div class="sidebar-kpi-title">Visual Extraction</div>
        <div class="sidebar-kpi-val" style="color:#00FFB2;">Images, Palette, Layout & Data</div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("---")
    st.markdown("### 📋 How It Works")
    st.markdown("""
    <div style="font-size:0.82rem; color:#94A3B8; line-height:1.6;">
        <b style="color:#FFF;">1. Ingest:</b> Upload brochures, specs, diagrams, or PDFs.<br>
        <b style="color:#FFF;">2. Analyze:</b> AI reads the content and decides its own color palette, fonts, icons, and per-slide layout.<br>
        <b style="color:#FFF;">3. Edit & Download:</b> Review, swap images/layouts, tweak slides, then export .PPTX.
    </div>
    """, unsafe_allow_html=True)

# =========================================================
# 7. Main Hero Section
# =========================================================
st.markdown('<div class="grid-backdrop">', unsafe_allow_html=True)
st.markdown('<div class="badge-pill"><span class="dot"></span>AI MULTIMODAL SYNTHESIZER</div>', unsafe_allow_html=True)
st.markdown('<div class="hero-title">VisualDeck AI</div>', unsafe_allow_html=True)
st.markdown('<div class="hero-subtitle">Transform unstructured visual documents, diagrams, and reports into structured, native PowerPoint presentations — with a design system the AI chooses for itself.</div>', unsafe_allow_html=True)
st.markdown('</div>', unsafe_allow_html=True)

# =========================================================
# 8. File Upload Section
# =========================================================
uploaded_files = st.file_uploader(
    label="Upload Visual Documents or Data Sheets",
    type=["png", "jpg", "jpeg", "pdf"],
    accept_multiple_files=True
)

image_count = 0
pdf_page_count = 0
total_size_mb = 0.0
raw_image_parts = []
source_images = []  # PIL images, index-aligned with raw_image_parts -> used later for embedding
original_files = []

if uploaded_files:
    total_size_mb = sum([f.size for f in uploaded_files]) / (1024 * 1024)
    for f in uploaded_files:
        if f.type.startswith("image/"):
            image_count += 1
            img = Image.open(f)
            original_files.append({"type": "image", "name": f.name, "size": f.size, "obj": img})

            buf = io.BytesIO()
            img_rgb = img.convert('RGB')
            img_rgb.save(buf, format='JPEG', quality=90)
            part = types.Part.from_bytes(data=buf.getvalue(), mime_type="image/jpeg")
            raw_image_parts.append(part)
            source_images.append(img_rgb)

        elif f.type == "application/pdf":
            try:
                reader = PdfReader(io.BytesIO(f.getvalue()))
                pages = len(reader.pages)
            except Exception:
                pages = 1
            pdf_page_count += pages
            original_files.append({"type": "pdf", "name": f.name, "size": f.size, "pages": pages})

# =========================================================
# 9. Top KPI Summary Cards
# =========================================================
col1, col2, col3, col4 = st.columns(4)
with col1:
    st.markdown("""<div class="kpi-card"><div class="kpi-title">AI Engine</div><div class="kpi-value">Gemini 3 Series</div><div class="kpi-desc">Multimodal Vision Pipeline</div></div>""", unsafe_allow_html=True)
with col2:
    val_str = f"{len(uploaded_files)} Items" if uploaded_files else "0 Files"
    desc_str = f"{image_count} Img | {pdf_page_count} PDF Pages" if pdf_page_count > 0 else f"{image_count} Images Loaded"
    st.markdown(f"""<div class="kpi-card"><div class="kpi-title">Input Buffer</div><div class="kpi-value">{val_str}</div><div class="kpi-desc">{desc_str}</div></div>""", unsafe_allow_html=True)
with col3:
    st.markdown(f"""<div class="kpi-card"><div class="kpi-title">Buffer Size</div><div class="kpi-value">{total_size_mb:.2f} MB</div><div class="kpi-desc">Ready for Synthesis</div></div>""", unsafe_allow_html=True)
with col4:
    st.markdown("""<div class="kpi-card"><div class="kpi-title">Export Standard</div><div class="kpi-value">16:9 .PPTX</div><div class="kpi-desc">Editable Shapes &amp; Images</div></div>""", unsafe_allow_html=True)

st.markdown("---")

# =========================================================
# 10. Source Assets Compact List
# =========================================================
if uploaded_files:
    st.markdown(f"### 📄 Uploaded Files ({len(uploaded_files)})")

    for idx, item in enumerate(original_files):
        with st.container(border=True):
            fcol1, fcol2 = st.columns([5, 1])
            with fcol1:
                size_kb = round(item["size"] / 1024, 1)
                icon = "🖼️" if item["type"] == "image" else "📄"
                st.markdown(f"**{icon} {item['name']}**")
                st.caption(f"{size_kb} KB" + (f" • {item['pages']} Pages" if item["type"] == "pdf" else ""))
            with fcol2:
                if item["type"] == "image":
                    if st.button("👁️ View", key=f"view_btn_{idx}", use_container_width=True):
                        preview_modal(item["obj"], item["name"])
                else:
                    st.button("📄 PDF", key=f"view_btn_{idx}", disabled=True, use_container_width=True)

    st.markdown("---")

    # =========================================================
    # 11. Presentation Customization Console
    # =========================================================
    st.markdown("### ⚙️ Presentation Customization")
    with st.container(border=True):
        opt_col1, opt_col2, opt_col3 = st.columns(3)
        with opt_col1:
            st.markdown("**🎯 Step 1: Slide Count**")
            slide_count = st.slider("Target number of slides:", min_value=3, max_value=14, value=6)
        with opt_col2:
            st.markdown("**🗣️ Step 2: Audience Tone**")
            presentation_tone = st.selectbox(
                "Presentation tone:",
                ["Executive Summary", "Product Deep Dive", "Investor Pitch", "Strategic Overview"],
                index=0
            )
        with opt_col3:
            st.markdown("**🎨 Step 3: Visual Style**")
            style_hint = st.selectbox(
                "Let the AI choose the palette based on:",
                ["Auto-Detect From Uploaded Visuals", "Bold & Futuristic", "Elegant & Minimal", "Corporate & Trustworthy", "Warm & Approachable"],
                index=0
            )

    generate_clicked = st.button("⚡ Generate AI Presentation Deck", type="primary")

    if generate_clicked:
        api_key = st.secrets.get("GEMINI_API_KEY")
        if not api_key:
            st.error("⚠️ Gemini API Key not found in Streamlit Secrets. Please configure `GEMINI_API_KEY`.")
        elif not raw_image_parts:
            st.warning("Please upload at least one image to synthesize slides.")
        else:
            loader_placeholder = st.empty()
            render_loader(loader_placeholder, "Synthesizing your deck", "PHASE 1 · Reading visuals, tables & brand cues...")

            client = genai.Client(api_key=api_key)

            style_instruction = (
                "Infer the palette, mood, and typography purely from what you see in the uploaded visuals "
                "(product photography, logos, existing brand colors, subject matter)."
                if style_hint == "Auto-Detect From Uploaded Visuals"
                else f"Lean toward a '{style_hint}' visual direction while still grounding choices in what you see in the uploaded visuals."
            )

            prompt = f"""
            You are an award-winning presentation designer and brand strategist. Analyze the uploaded source
            images carefully — their subject matter, mood, existing colors/logos, and level of formality.

            Your job is to make ALL design decisions yourself. Nothing about color, font, icon, or layout choice
            should be generic or copy-pasted between decks — base every choice on the actual content you see.

            GUIDELINES:
            1. FLOW: Produce a {slide_count}-slide outline matching this tone: '{presentation_tone}'.
            2. DATA EXTRACTION: Extract exact model names, capacities, certifications, specs, and pricing into
               native tables and KPI metrics. Never invent numbers that are not visible or reasonably implied.
            3. COLOR & TYPE: {style_instruction} Choose background_hex, primary_accent_hex, secondary_accent_hex,
               card_bg_hex, text_color_hex, and secondary_text_hex so they are harmonious and have strong contrast
               for readability on a dark or light canvas (your choice). Choose heading_font and body_font ONLY from
               this exact list (PowerPoint-safe fonts): Calibri, Arial, Segoe UI, Georgia, Verdana, Trebuchet MS,
               Century Gothic, Garamond, Tahoma — pick a pairing that matches the mood (e.g. Georgia for elegant,
               Century Gothic for modern/tech, Calibri for corporate).
            4. ICONS: For every KPI and bullet, choose ONE short unicode symbol or emoji in the "icon" field that
               best represents that specific point (not a generic bullet) — pick freely based on meaning.
            5. IMAGES: You are given {len(raw_image_parts)} source image(s), indexed 0 to {len(raw_image_parts) - 1}
               in the order provided. For any slide where showing one of these images would help
               (e.g. showcasing a specific product), set "layout_type" to "image_left", "image_right", or
               "image_full" and set "image_index" to the matching image's index. Only reference indices that exist.
               Do not force an image onto every slide — use table/kpi_grid/cards layouts when that communicates
               the content better. You may also use "section_break" once for a clean divider slide if it improves
               pacing for a longer deck.
            6. LAYOUT VARIETY: Do not repeat the same layout_type on every slide unless the content truly calls
               for it. Choose whichever of: kpi_grid, table, cards, image_left, image_right, image_full,
               section_break best fits each slide's specific content.
            7. BULLETS FORMAT: Every bullet must be an object: {{"icon": "🔧", "text": "..."}} — never a bare string.

            Output ONLY a JSON object conforming strictly to this structure (no markdown fences, no commentary):
            {{
              "theme": {{
                "background_hex": "#0A0E17",
                "primary_accent_hex": "#00D2FF",
                "secondary_accent_hex": "#9D00FF",
                "card_bg_hex": "#141C2E",
                "text_color_hex": "#FFFFFF",
                "secondary_text_hex": "#94A3B8",
                "heading_font": "Century Gothic",
                "body_font": "Calibri"
              }},
              "deck_title": "Executive Presentation Title",
              "slides": [
                {{
                  "slide_number": 1,
                  "layout_type": "kpi_grid",
                  "title": "Slide Title",
                  "subtitle": "Subtitle or core insight",
                  "kpis": [
                    {{"label": "Metric Name", "value": "Extracted Value", "desc": "Context note", "icon": "⚙️"}}
                  ]
                }},
                {{
                  "slide_number": 2,
                  "layout_type": "image_left",
                  "title": "Slide Title",
                  "subtitle": "Subtitle",
                  "image_index": 0,
                  "bullets": [
                    {{"icon": "✅", "text": "Point about the product shown"}}
                  ]
                }},
                {{
                  "slide_number": 3,
                  "layout_type": "table",
                  "title": "Model Comparison & Specifications",
                  "subtitle": "Detailed breakdown",
                  "table": {{
                    "headers": ["Model / Item", "Capacity / Spec", "Investment / Price"],
                    "rows": [["Model A", "2000L", "Rs 7,49,300/-"]]
                  }}
                }},
                {{
                  "slide_number": 4,
                  "layout_type": "cards",
                  "title": "Key Advantages & Features",
                  "subtitle": "Value summary",
                  "bullets": [
                    {{"icon": "🛡️", "text": "Clear bullet point 1"}}
                  ]
                }}
              ]
            }}
            """

            gemini_models = [
                'gemini-3.1-flash-lite',
                'gemini-3.5-flash-lite',
                'gemini-3.7-flash',
                'gemini-3.5-flash'
            ]

            response_json = None
            last_error = None

            for model_name in gemini_models:
                try:
                    content_payload = [prompt] + raw_image_parts[:5]
                    res = client.models.generate_content(
                        model=model_name,
                        contents=content_payload,
                        config=types.GenerateContentConfig(
                            response_mime_type="application/json"
                        )
                    )
                    raw_text = res.text.strip()
                    response_json = json.loads(raw_text)
                    break
                except Exception as e:
                    last_error = str(e)
                    continue

            if not response_json:
                loader_placeholder.empty()
                st.error(f"API Error Log: {last_error}")
            else:
                render_loader(loader_placeholder, "Synthesizing your deck", "PHASE 2 · Building native slides, images & tables...")
                pptx_stream = build_powerpoint(response_json, source_images)
                st.session_state["generated_pptx"] = pptx_stream
                st.session_state["deck_data"] = response_json
                st.session_state["source_images"] = source_images
                time.sleep(0.4)
                loader_placeholder.empty()
                st.success("✅ Presentation deck synthesized successfully!")

    # =========================================================
    # 12. Interactive Slide Preview & Live Deck Editor
    # =========================================================
    if "deck_data" in st.session_state and st.session_state["deck_data"]:
        deck = st.session_state["deck_data"]
        slides = deck.get("slides", [])
        src_imgs = st.session_state.get("source_images", source_images)

        st.markdown("---")
        st.markdown(f"### 📑 Synthesized Presentation: **{deck.get('deck_title', 'Synthesized Deck')}**")

        st.download_button(
            label="📥 Download Native 16:9 .PPTX Presentation",
            data=st.session_state["generated_pptx"],
            file_name="VisualDeck_AI_Presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary"
        )

        st.markdown("#### ✏️ Interactive Slide Preview & Editor")
        st.caption("Review your generated slides below. Switch layouts, swap images, or edit text/tables before downloading.")

        LAYOUT_OPTIONS = ["kpi_grid", "table", "cards", "image_left", "image_right", "image_full", "section_break"]

        for idx, s in enumerate(slides):
            with st.container(border=True):
                st.markdown(f"### 🖥️ Slide {s.get('slide_number', idx + 1)}: {s.get('title', '')}")

                with st.expander("👁️ View Rendered Slide Card", expanded=True):
                    st.markdown(f"**{s.get('title', '')}**")
                    if s.get("subtitle"):
                        st.caption(f"_{s.get('subtitle')}_")

                    if s.get("layout_type") in ("image_left", "image_right", "image_full") and isinstance(s.get("image_index"), int) and src_imgs:
                        i = s["image_index"]
                        if 0 <= i < len(src_imgs):
                            st.image(src_imgs[i], width=320)

                    if s.get("layout_type") == "kpi_grid" and s.get("kpis"):
                        kpi_cols = st.columns(len(s["kpis"]))
                        for k_idx, k in enumerate(s["kpis"]):
                            with kpi_cols[k_idx]:
                                st.metric(label=f"{k.get('icon', '')} {k.get('label', 'Metric')}", value=k.get("value", "N/A"))

                    elif s.get("layout_type") == "table" and s.get("table"):
                        headers = s["table"].get("headers", [])
                        rows = s["table"].get("rows", [])
                        if headers and rows:
                            df = pd.DataFrame(rows, columns=headers)
                            st.dataframe(df, use_container_width=True)

                    elif s.get("bullets"):
                        for b in s["bullets"]:
                            if isinstance(b, dict):
                                st.markdown(f"{b.get('icon', '•')} {b.get('text', '')}")
                            else:
                                st.markdown(f"• {b}")

                # ---- Easier, tabbed editing controls ----
                with st.expander(f"✏️ Edit Slide {idx + 1}"):
                    tab_content, tab_layout = st.tabs(["📝 Content", "🎛️ Layout & Image"])

                    with tab_content:
                        s["title"] = st.text_input("Slide Title:", value=s.get("title", ""), key=f"title_{idx}")
                        s["subtitle"] = st.text_input("Subtitle / Insight:", value=s.get("subtitle", ""), key=f"sub_{idx}")

                        if s.get("layout_type") == "table" and s.get("table"):
                            st.markdown("**Edit Table Data:**")
                            headers = s["table"].get("headers", [])
                            rows = s["table"].get("rows", [])
                            if headers and rows:
                                df_edit = pd.DataFrame(rows, columns=headers)
                                edited_df = st.data_editor(df_edit, key=f"tbl_edit_{idx}", num_rows="dynamic")
                                s["table"]["headers"] = list(edited_df.columns)
                                s["table"]["rows"] = edited_df.values.tolist()

                        elif s.get("layout_type") in ("cards", "image_left", "image_right") and s.get("bullets"):
                            st.markdown("**Edit Points (one per line — optionally start with an emoji icon):**")
                            bullet_lines = []
                            for b in s.get("bullets", []):
                                if isinstance(b, dict):
                                    bullet_lines.append(f"{b.get('icon', '')} {b.get('text', '')}".strip())
                                else:
                                    bullet_lines.append(str(b))
                            edited_text = st.text_area("Points:", value="\n".join(bullet_lines), key=f"bullets_{idx}", height=120)
                            new_bullets = []
                            for line in edited_text.split("\n"):
                                line = line.strip()
                                if not line:
                                    continue
                                parts = line.split(" ", 1)
                                if len(parts) == 2 and len(parts[0]) <= 4:
                                    new_bullets.append({"icon": parts[0], "text": parts[1]})
                                else:
                                    new_bullets.append({"icon": "•", "text": line})
                            s["bullets"] = new_bullets

                        elif s.get("layout_type") == "kpi_grid" and s.get("kpis"):
                            st.markdown("**Edit KPI Cards:**")
                            for k_idx, k in enumerate(s["kpis"]):
                                kcols = st.columns([1, 2, 2, 3])
                                k["icon"] = kcols[0].text_input("Icon", value=k.get("icon", ""), key=f"kpi_icon_{idx}_{k_idx}")
                                k["label"] = kcols[1].text_input("Label", value=k.get("label", ""), key=f"kpi_lbl_{idx}_{k_idx}")
                                k["value"] = kcols[2].text_input("Value", value=str(k.get("value", "")), key=f"kpi_val_{idx}_{k_idx}")
                                k["desc"] = kcols[3].text_input("Description", value=k.get("desc", ""), key=f"kpi_desc_{idx}_{k_idx}")

                    with tab_layout:
                        current_layout = s.get("layout_type", "cards")
                        new_layout = st.selectbox(
                            "Slide layout:", LAYOUT_OPTIONS,
                            index=LAYOUT_OPTIONS.index(current_layout) if current_layout in LAYOUT_OPTIONS else 0,
                            key=f"layout_{idx}"
                        )
                        s["layout_type"] = new_layout

                        if new_layout in ("image_left", "image_right", "image_full") and src_imgs:
                            img_labels = [f"Image {i}" for i in range(len(src_imgs))]
                            current_idx = s.get("image_index", 0) if isinstance(s.get("image_index"), int) else 0
                            current_idx = current_idx if 0 <= current_idx < len(src_imgs) else 0
                            chosen = st.selectbox("Image to use:", img_labels, index=current_idx, key=f"img_pick_{idx}")
                            s["image_index"] = img_labels.index(chosen)
                            st.image(src_imgs[s["image_index"]], width=220)

        if st.button("💾 Apply Edits & Rebuild Presentation", type="secondary"):
            st.session_state["generated_pptx"] = build_powerpoint(st.session_state["deck_data"], src_imgs)
            st.success("✅ Presentation updated with your changes! Click Download above to get your updated .pptx file.")
            st.rerun()

else:
    st.info("👆 Upload one or more product sheets, diagrams, or PDF documents above to begin.")
